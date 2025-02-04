import os
import pandas as pd
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt

# Create output directory if it doesn't exist
output_dir = "sample_docs"
os.makedirs(output_dir, exist_ok=True)

#################################
# 1. Create patient_records_sample.xlsx with more rows
#################################
patient_data = {
    "PatientID": [1001, 1002, 1003, 1004, 1005],
    "Name": ["John Doe", "Jane Smith", "Alice Johnson", "Bob Brown", "Carol White"],
    "AdmissionDate": ["2024-03-15", "2024-03-20", "2024-04-01", "2024-04-05", "2024-04-10"],
    "Diagnosis": [
        "Hypertension", 
        "Type 2 Diabetes", 
        "Asthma", 
        "Coronary Artery Disease", 
        "Chronic Kidney Disease"
    ],
    "Treatment": [
        "Lisinopril 10mg daily", 
        "Metformin 500mg BID", 
        "Albuterol inhaler as needed", 
        "Aspirin 81mg daily, beta-blockers", 
        "ACE inhibitors and dietary modifications"
    ]
}
df_patients = pd.DataFrame(patient_data)
patient_excel_path = os.path.join(output_dir, "patient_records_sample.xlsx")
df_patients.to_excel(patient_excel_path, index=False)
print(f"Created {patient_excel_path}")

#################################
# 2. Create medication_inventory.xlsx with additional rows
#################################
medication_data = {
    "DrugID": ["M101", "M202", "M303", "M404"],
    "DrugName": ["Lisinopril", "Metformin", "Atorvastatin", "Amlodipine"],
    "Quantity": [1500, 3000, 2000, 1800],
    "ExpiryDate": ["2025-12-31", "2026-06-30", "2025-09-30", "2026-03-31"],
    "Supplier": ["MedSupply Co.", "PharmaDistro Inc.", "HealthMeds Ltd.", "Global Pharma"]
}
df_med = pd.DataFrame(medication_data)
med_excel_path = os.path.join(output_dir, "medication_inventory.xlsx")
df_med.to_excel(med_excel_path, index=False)
print(f"Created {med_excel_path}")

#################################
# 3. Create clinical_guidelines.pdf with more content
#################################
clinical_guidelines_text = (
    "Hypertension Management Guidelines\n\n"
    "1. Target Blood Pressure:\n"
    "   - <130/80 mmHg for most patients\n"
    "   - <140/90 mmHg for elderly patients\n\n"
    "2. First-line Medications:\n"
    "   1. ACE Inhibitors\n"
    "   2. Calcium Channel Blockers\n"
    "   3. Thiazide Diuretics\n\n"
    "3. Lifestyle Modifications:\n"
    "   - Weight loss\n"
    "   - Low-sodium diet\n"
    "   - Regular physical activity\n\n"
    "4. Follow-Up and Monitoring:\n"
    "   - Regular blood pressure monitoring\n"
    "   - Adjust medications based on patient response\n"
)
pdf_path_guidelines = os.path.join(output_dir, "clinical_guidelines.pdf")
c = canvas.Canvas(pdf_path_guidelines, pagesize=LETTER)
width, height = LETTER

# Draw title
c.setFont("Helvetica-Bold", 16)
c.drawString(1 * inch, height - 1 * inch, "Hypertension Management Guidelines")
c.setFont("Helvetica", 12)
text_obj = c.beginText(1 * inch, height - 1.5 * inch)
for line in clinical_guidelines_text.splitlines():
    text_obj.textLine(line)
c.drawText(text_obj)
c.showPage()
c.save()
print(f"Created {pdf_path_guidelines}")

#################################
# 4. Create vendor_contract_agreement.pdf with additional content
#################################
vendor_contract_text = (
    "Supplier Agreement - MedSupply Co.\n\n"
    "Term: 2024-2026\n\n"
    "Delivery SLAs:\n"
    "   - Emergency Orders: 24hr delivery\n"
    "   - Regular Orders: 5 business days\n\n"
    "Pricing Terms:\n"
    "   - 2% discount for payments within 15 days\n\n"
    "Additional Clauses:\n"
    "   - Confidentiality: Both parties agree to keep pricing and contract details confidential.\n"
    "   - Performance Metrics: Monthly performance reviews will be held.\n"
    "   - Termination: Either party may terminate the agreement with a 30-day notice if performance standards are not met.\n"
)
pdf_path_contract = os.path.join(output_dir, "vendor_contract_agreement.pdf")
c = canvas.Canvas(pdf_path_contract, pagesize=LETTER)
c.setFont("Helvetica-Bold", 16)
c.drawString(1 * inch, height - 1 * inch, "Supplier Agreement - MedSupply Co.")
c.setFont("Helvetica", 12)
text_obj = c.beginText(1 * inch, height - 1.5 * inch)
for line in vendor_contract_text.splitlines():
    text_obj.textLine(line)
c.drawText(text_obj)
c.showPage()
c.save()
print(f"Created {pdf_path_contract}")

#################################
# 5. Create infection_control_training.pptx with more slides
#################################
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hospital Infection Control"
subtitle.text = "Infection Control Training for Healthcare Providers"

# Slide 2: Key Protocols
slide_layout = prs.slide_layouts[1]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Key Protocols"
content.text = "• Hand hygiene compliance >95%\n• PPE requirements\n• Environmental cleaning schedules"

# Slide 3: Outbreak Management
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Outbreak Management"
content.text = "• Isolation procedures\n• Rapid response teams\n• Communication protocols"

# Slide 4: Staff Training and Drills
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Staff Training and Drills"
content.text = "• Regular training sessions\n• Simulation drills\n• Evaluation and feedback"

# Slide 5: Waste Management and Disposal
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Waste Management"
content.text = "• Proper disposal of medical waste\n• Recycling protocols\n• Safety measures"

# Slide 6: Communication and Reporting
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Communication & Reporting"
content.text = "• Incident reporting procedures\n• Daily safety briefings\n• Coordination with local health authorities"

# Slide 7: Conclusion and Q&A
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.shapes.placeholders[1]
title.text = "Conclusion & Q&A"
content.text = "• Review of key protocols\n• Open discussion\n• Next steps for improvement"

ppt_path = os.path.join(output_dir, "infection_control_training.pptx")
prs.save(ppt_path)
print(f"Created {ppt_path}")